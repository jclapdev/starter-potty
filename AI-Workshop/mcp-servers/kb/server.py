#!/usr/bin/env python3
"""KB MCP server — queryable knowledge base of tool problems and fixes.

A shared knowledge layer: agents look up how a past problem with a specific
tool was solved, and record new fixes as they resolve them. Backed by a local
vector store (LanceDB) with local embeddings (sentence-transformers), so there
is no API cost and nothing leaves the machine.

Tools (5):
  read
    - query_patterns(tool, problem, limit=5)  → ranked problem/fix matches
    - list_tools()                            → tools present + record counts
    - get_patterns_summary(tool)              → count, top tags, date range
  write
    - record_fix(tool, problem, fix, ...)     → store a problem/fix (dedup'd)
    - ingest_kb(source_path, tool=None)       → ingest a file or directory

Auto-ingestion: on startup the server scans KB_RESOURCES_DIR for new/changed
.md/.txt/.pdf files and ingests them before serving. No file watcher — a
restart re-scans. Failures during the scan are logged to stderr and never
abort startup.

The MCP stdio loop (_send / _serve_stdio) is the same newline-delimited
JSON-RPC 2.0 protocol used by the vault-mcp server.

Heavy deps (lancedb, sentence-transformers, pypdf) are imported lazily so the
process starts fast and tools that do not touch the store stay cheap.

Paths come from env vars so each machine configures its own:
  KB_DATA_PATH       LanceDB storage      (default: <server_dir>/data)
  KB_RESOURCES_PATH  auto-ingest folder   (default: <server_dir>/sources)
  KB_MODEL_NAME      embedding model      (default: all-MiniLM-L6-v2)
"""
from __future__ import annotations

import hashlib
import json
import os
import re
import sys
from collections import Counter
from datetime import datetime, timezone
from pathlib import Path

_SERVER_DIR = Path(__file__).resolve().parent
KB_DATA_DIR = Path(os.environ.get("KB_DATA_PATH", str(_SERVER_DIR / "data")))
KB_RESOURCES_DIR = Path(os.environ.get("KB_RESOURCES_PATH", str(_SERVER_DIR / "sources")))
KB_MODEL_NAME = os.environ.get("KB_MODEL_NAME", "all-MiniLM-L6-v2")

# Vault root (server lives at <vault>/AI-Workshop/mcp-servers/kb/server.py).
KB_VAULT_DIR = Path(os.environ.get("KB_VAULT_PATH", str(_SERVER_DIR.parent.parent.parent)))
# Vault folders indexed on startup (tagged tool="vault") for semantic recall
# over the system's own history, decisions, and memory. Override with
# KB_INGEST_PATHS (os.pathsep-separated) — set to empty to disable.
_default_ingest = os.pathsep.join([str(KB_VAULT_DIR / "Context"), str(KB_VAULT_DIR / "Reference")])
KB_INGEST_PATHS = [Path(p) for p in
                   os.environ.get("KB_INGEST_PATHS", _default_ingest).split(os.pathsep)
                   if p.strip()]
# Where "index+note" transpose mode writes clean Obsidian notes (the wiki layer).
KB_WIKI_DIR = Path(os.environ.get("KB_WIKI_PATH",
                                  str(KB_VAULT_DIR / "AI-Workshop" / "Projects" / "Wiki")))

_VECTOR_DIM = 384  # all-MiniLM-L6-v2 output width
_SCAN_LIMIT = 1_000_000_000  # effectively "all rows" for column-projected scans
_RECORDS_TABLE = "kb_records"
_INGEST_TABLE = "ingested_files"

# Lazy-loaded singletons — populated on first use, not at import.
_MODEL = None
_DB = None
_TABLE = None
_ING = None


# --------------------------------------------------------------------------- #
# Small helpers
# --------------------------------------------------------------------------- #
def _log(msg):
    """Diagnostics go to stderr so they never corrupt the JSON-RPC stream."""
    sys.stderr.write("[kb-mcp] %s\n" % msg)
    sys.stderr.flush()


def _now():
    return datetime.now(timezone.utc).isoformat()


def _esc(s):
    """Escape a value for a LanceDB SQL filter literal (single quotes doubled)."""
    return str(s).replace("'", "''")


def _hash(text):
    return hashlib.sha256(text.encode("utf-8")).hexdigest()


def _loadtags(raw):
    try:
        v = json.loads(raw or "[]")
        return v if isinstance(v, list) else []
    except (ValueError, TypeError):
        return []


# --------------------------------------------------------------------------- #
# Lazy resources: embedding model, DB, tables
# --------------------------------------------------------------------------- #
def _schema():
    """PyArrow schema for kb_records. Defined here so pyarrow loads lazily."""
    import pyarrow as pa
    return pa.schema([
        ("id", pa.string()),
        ("vector", pa.list_(pa.float32(), _VECTOR_DIM)),
        ("tool", pa.string()),
        ("project", pa.string()),
        ("type", pa.string()),
        ("problem", pa.string()),
        ("fix", pa.string()),
        ("tags", pa.string()),
        ("source_path", pa.string()),
        ("date_added", pa.string()),
    ])


def _ingest_schema():
    import pyarrow as pa
    return pa.schema([
        ("path", pa.string()),
        ("mtime", pa.float64()),
        ("chunk_count", pa.int64()),
        ("date_ingested", pa.string()),
    ])


def _get_model():
    global _MODEL
    if _MODEL is None:
        from sentence_transformers import SentenceTransformer
        _log("loading embedding model %s" % KB_MODEL_NAME)
        _MODEL = SentenceTransformer(KB_MODEL_NAME)
    return _MODEL


def _embed(text):
    """Return a normalized 384-d embedding as a plain python list."""
    return _get_model().encode(text, normalize_embeddings=True).tolist()


def _connect():
    global _DB
    if _DB is None:
        import lancedb
        KB_DATA_DIR.mkdir(parents=True, exist_ok=True)
        _DB = lancedb.connect(str(KB_DATA_DIR))
    return _DB


def _get_table():
    """kb_records — created on first use (first-run safe)."""
    global _TABLE
    if _TABLE is None:
        db = _connect()
        if _RECORDS_TABLE in db.table_names():
            _TABLE = db.open_table(_RECORDS_TABLE)
        else:
            _TABLE = db.create_table(_RECORDS_TABLE, schema=_schema())
    return _TABLE


def _get_ingest_table():
    global _ING
    if _ING is None:
        db = _connect()
        if _INGEST_TABLE in db.table_names():
            _ING = db.open_table(_INGEST_TABLE)
        else:
            _ING = db.create_table(_INGEST_TABLE, schema=_ingest_schema())
    return _ING


# --------------------------------------------------------------------------- #
# Text extraction + chunking
# --------------------------------------------------------------------------- #
# --- Extractor registry ---------------------------------------------------- #
# Each extractor maps a file extension to a function returning plain text (or,
# for tables, a list of pre-formed row-chunks). Optional libraries are imported
# lazily *inside* each extractor: a type whose library isn't installed is logged
# and skipped, never crashing ingestion. Pure-Python libraries only — no
# shelling out to LibreOffice / pandoc / Tesseract, so the setup stays portable.

def _x_plain(path):
    return path.read_text(encoding="utf-8", errors="replace")


def _x_pdf(path):
    from pypdf import PdfReader
    reader = PdfReader(str(path))
    return "\n\n".join((pg.extract_text() or "") for pg in reader.pages)


def _x_docx(path):
    import docx  # python-docx
    doc = docx.Document(str(path))
    parts = [p.text for p in doc.paragraphs if p.text.strip()]
    for table in doc.tables:
        for row in table.rows:
            cells = [c.text.strip() for c in row.cells if c.text.strip()]
            if cells:
                parts.append(" | ".join(cells))
    return "\n\n".join(parts)


def _x_pptx(path):
    from pptx import Presentation  # python-pptx
    prs = Presentation(str(path))
    parts = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame and shape.text_frame.text.strip():
                parts.append(shape.text_frame.text.strip())
    return "\n\n".join(parts)


def _x_html(path):
    from bs4 import BeautifulSoup
    soup = BeautifulSoup(path.read_text(encoding="utf-8", errors="replace"), "html.parser")
    for tag in soup(["script", "style", "noscript"]):
        tag.decompose()
    return soup.get_text("\n")


def _x_rtf(path):
    from striprtf.striprtf import rtf_to_text
    return rtf_to_text(path.read_text(encoding="utf-8", errors="replace"))


def _x_epub(path):
    from ebooklib import epub, ITEM_DOCUMENT
    from bs4 import BeautifulSoup
    book = epub.read_epub(str(path))
    parts = []
    for item in book.get_items_of_type(ITEM_DOCUMENT):
        soup = BeautifulSoup(item.get_content(), "html.parser")
        txt = soup.get_text("\n").strip()
        if txt:
            parts.append(txt)
    return "\n\n".join(parts)


def _x_json(path):
    import json as _json
    data = _json.loads(path.read_text(encoding="utf-8", errors="replace"))
    lines = []

    def walk(obj, prefix):
        if isinstance(obj, dict):
            for k, v in obj.items():
                walk(v, "%s.%s" % (prefix, k) if prefix else str(k))
        elif isinstance(obj, list):
            for i, v in enumerate(obj):
                walk(v, "%s[%d]" % (prefix, i))
        else:
            lines.append("%s: %s" % (prefix, obj) if prefix else str(obj))

    walk(data, "")
    return "\n".join(lines)


def _x_xml(path):
    import xml.etree.ElementTree as ET
    root = ET.parse(str(path)).getroot()
    lines = []

    def walk(elem, prefix):
        tag = elem.tag.split("}")[-1]  # strip namespace
        here = "%s/%s" % (prefix, tag) if prefix else tag
        if elem.text and elem.text.strip():
            lines.append("%s: %s" % (here, elem.text.strip()))
        for k, v in elem.attrib.items():
            lines.append("%s@%s: %s" % (here, k.split("}")[-1], v))
        for child in elem:
            walk(child, here)

    walk(root, "")
    return "\n".join(lines)


def _is_number(val):
    if isinstance(val, (int, float)):
        return True
    try:
        float(str(val).replace(",", "").strip())
        return True
    except (ValueError, AttributeError):
        return False


def _rows_to_chunks(header, data_rows, sheet=None):
    """One chunk per row. Keep only text cells (numeric-only tables are filter
    queries, the wrong tool for a vector store); drop a row with no text."""
    prefix = ("[%s] " % sheet) if sheet else ""
    chunks = []
    for row in data_rows:
        cells = []
        for i, val in enumerate(row):
            if val is None or str(val).strip() == "" or _is_number(val):
                continue
            col = str(header[i]) if i < len(header) and header[i] not in (None, "") else "col%d" % i
            cells.append("%s: %s" % (col, str(val).strip()))
        if cells:
            chunks.append(prefix + " | ".join(cells))
    return chunks


def _x_csv(path):
    import csv
    with path.open(encoding="utf-8", errors="replace", newline="") as fh:
        rows = list(csv.reader(fh))
    if not rows:
        return []
    return _rows_to_chunks(rows[0], rows[1:])


def _x_xlsx(path):
    from openpyxl import load_workbook
    wb = load_workbook(str(path), read_only=True, data_only=True)
    chunks = []
    for ws in wb.worksheets:
        rows = [list(r) for r in ws.iter_rows(values_only=True)]
        if not rows:
            continue
        chunks.extend(_rows_to_chunks(rows[0], rows[1:], sheet=ws.title))
    wb.close()
    return chunks


_OCR_ENGINE = None


def _get_ocr():
    global _OCR_ENGINE
    if _OCR_ENGINE is None:
        from rapidocr_onnxruntime import RapidOCR
        _OCR_ENGINE = RapidOCR()
    return _OCR_ENGINE


def _x_image(path):
    result, _ = _get_ocr()(str(path))
    if not result:
        return ""
    return "\n".join(line[1] for line in result)


# extension -> extractor returning a single text blob (chunked downstream)
_BLOB_EXTRACTORS = {
    ".md": _x_plain, ".txt": _x_plain, ".pdf": _x_pdf,
    ".docx": _x_docx, ".pptx": _x_pptx,
    ".html": _x_html, ".htm": _x_html,
    ".rtf": _x_rtf, ".epub": _x_epub,
    ".json": _x_json, ".xml": _x_xml,
}
# extension -> extractor returning a list of pre-formed row-chunks (used as-is)
_ROW_EXTRACTORS = {".csv": _x_csv, ".xlsx": _x_xlsx}
# OCR'd as a single text blob (chunked downstream)
_IMAGE_EXT = {".png", ".jpg", ".jpeg", ".webp", ".bmp", ".tiff", ".tif"}

_SUPPORTED_EXT = set(_BLOB_EXTRACTORS) | set(_ROW_EXTRACTORS) | _IMAGE_EXT


def _extract_chunks(path):
    """Return a list of ready-to-embed text chunks for a file. Returns [] for an
    unsupported type, a missing optional library, or an extraction error — none
    of which abort ingestion (each is logged to stderr)."""
    suf = path.suffix.lower()
    try:
        if suf in _ROW_EXTRACTORS:
            return [c for c in _ROW_EXTRACTORS[suf](path) if c.strip()]
        if suf in _IMAGE_EXT:
            text = _x_image(path)
            return _chunk_text(text) if text.strip() else []
        if suf in _BLOB_EXTRACTORS:
            text = _BLOB_EXTRACTORS[suf](path)
            return _chunk_text(text) if text.strip() else []
    except ImportError as exc:
        _log("skip %s: missing library for %s (%s)" % (path.name, suf, exc))
        return []
    except Exception as exc:  # noqa: BLE001 — a bad file must never abort a batch
        _log("skip %s: extract error (%s)" % (path.name, exc))
        return []
    _log("skip %s: unsupported type %s" % (path.name, suf))
    return []


def _chunk_text(text, target=512, overlap=64):
    """Split text into ~target-word chunks on paragraph boundaries, carrying an
    overlap of the trailing words into the next chunk so context spans breaks.
    Word count is a cheap stand-in for token count — close enough for chunking.
    """
    paras = [p.strip() for p in re.split(r"\n\s*\n", text) if p.strip()]
    chunks = []
    cur = []
    cur_len = 0
    for p in paras:
        plen = len(p.split())
        if cur and cur_len + plen > target:
            chunk = "\n\n".join(cur)
            chunks.append(chunk)
            tail = chunk.split()[-overlap:] if overlap else []
            cur = [" ".join(tail)] if tail else []
            cur_len = len(tail)
        cur.append(p)
        cur_len += plen
    if cur:
        chunks.append("\n\n".join(cur))
    return [c for c in chunks if c.strip()]


def _chunk_heading(chunk):
    """First markdown heading or first non-empty line, as the chunk's `problem`."""
    for line in chunk.splitlines():
        s = line.strip().lstrip("#").strip()
        if s:
            return s[:200]
    return chunk[:200]


# --------------------------------------------------------------------------- #
# Ingestion pipeline (shared by ingest.py CLI and the ingest_kb tool)
# --------------------------------------------------------------------------- #
def _gather_files(path):
    path = Path(path)
    if path.is_file():
        return [path] if path.suffix.lower() in _SUPPORTED_EXT else []
    files = []
    for p in sorted(path.rglob("*")):
        if p.is_file() and p.suffix.lower() in _SUPPORTED_EXT:
            files.append(p)
    return files


def _slugify(text):
    return re.sub(r"[^a-z0-9]+", "-", text.lower()).strip("-") or "source"


def _transpose_to_note(src, chunks):
    """'index+note' mode: write the extracted text of a non-Markdown source into
    a clean Obsidian note under the wiki folder. Markdown/text sources are left
    alone (already Obsidian-native). Returns the note path, or None."""
    if src.suffix.lower() in (".md",):
        return None
    try:
        KB_WIKI_DIR.mkdir(parents=True, exist_ok=True)
        note = KB_WIKI_DIR / (_slugify(src.stem) + ".md")
        front = ("---\ntype: transposed-source\nsource: %s\ningested: %s\n"
                 "tags: [kb, transposed]\n---\n\n" % (src.name, _now()[:10]))
        note.write_text(front + "# %s\n\n%s\n" % (src.stem, "\n\n".join(chunks)),
                        encoding="utf-8")
        return note
    except Exception as exc:  # noqa: BLE001 — a note-write failure must not fail ingest
        _log("transpose note failed for %s: %s" % (src.name, exc))
        return None


def _process_file(path, tool=None, write_note=False):
    """Ingest a single file. Idempotent: unchanged files and duplicate chunks
    are skipped. With write_note, also transpose the source to an Obsidian note
    (index+note mode). Returns a per-file summary dict.
    """
    path = Path(path)
    try:
        mtime = path.stat().st_mtime
    except OSError as exc:
        return {"file": str(path), "ingested": 0, "skipped": "stat-error: %s" % exc}

    ing = _get_ingest_table()
    if ing.count_rows(filter="path = '%s' AND mtime = %r" % (_esc(str(path)), mtime)):
        return {"file": str(path), "ingested": 0, "skipped": "unchanged"}

    chunks = _extract_chunks(path)
    if not chunks:
        return {"file": str(path), "ingested": 0, "skipped": "empty-or-unsupported"}

    note_path = _transpose_to_note(path, chunks) if write_note else None

    tbl = _get_table()
    rows = []
    duplicates = 0
    derived_tool = tool or path.stem
    for ch in chunks:
        cid = _hash(ch)
        if tbl.count_rows(filter="id = '%s'" % cid):
            duplicates += 1
            continue
        rows.append({
            "id": cid,
            "vector": _embed(ch),
            "tool": derived_tool,
            "project": "",
            "type": "reference",
            "problem": _chunk_heading(ch),
            "fix": "",
            "tags": "[]",
            "source_path": str(path),
            "date_added": _now(),
        })
    if rows:
        tbl.add(rows)

    # Upsert the ingestion checkpoint: drop any prior row for this path, re-add.
    ing.delete("path = '%s'" % _esc(str(path)))
    ing.add([{
        "path": str(path),
        "mtime": mtime,
        "chunk_count": len(rows),
        "date_ingested": _now(),
    }])
    out = {"file": str(path), "ingested": len(rows), "duplicates": duplicates}
    if note_path:
        out["note"] = str(note_path)
    return out


# --------------------------------------------------------------------------- #
# Tool cores
# --------------------------------------------------------------------------- #
def list_tools_core():
    tbl = _get_table()
    if tbl.count_rows() == 0:
        return {"tools": []}
    # Project only the `tool` column so the vector column is never read.
    arrow = tbl.search().select(["tool"]).limit(_SCAN_LIMIT).to_arrow()
    counts = Counter(arrow.column("tool").to_pylist())
    return {"tools": [{"tool": t, "records": n}
                      for t, n in sorted(counts.items(), key=lambda x: -x[1])]}


def get_patterns_summary_core(tool):
    tbl = _get_table()
    flt = "tool = '%s'" % _esc(tool)
    n = tbl.count_rows(filter=flt)
    if n == 0:
        return {"tool": tool, "records": 0, "top_tags": [], "date_range": None}
    arrow = (tbl.search().where(flt).select(["tags", "date_added"])
             .limit(_SCAN_LIMIT).to_arrow())
    tags = Counter()
    for raw in arrow.column("tags").to_pylist():
        for t in _loadtags(raw):
            tags[t] += 1
    dates = [d for d in arrow.column("date_added").to_pylist() if d]
    return {
        "tool": tool,
        "records": n,
        "top_tags": [{"tag": k, "count": v} for k, v in tags.most_common(10)],
        "date_range": {"earliest": min(dates), "latest": max(dates)} if dates else None,
    }


def query_patterns_core(tool, problem, limit=5):
    tbl = _get_table()
    if tbl.count_rows() == 0:
        return {"tool": tool, "problem": problem, "matches": []}
    res = (tbl.search(_embed(problem))
           .metric("cosine")
           .where("tool = '%s'" % _esc(tool), prefilter=True)
           .limit(max(1, int(limit)))
           .to_list())
    matches = []
    for r in res:
        dist = r.get("_distance")
        matches.append({
            "problem": r.get("problem"),
            "fix": r.get("fix"),
            "type": r.get("type"),
            "project": r.get("project") or None,
            "tags": _loadtags(r.get("tags")),
            "source_path": r.get("source_path") or None,
            "date_added": r.get("date_added"),
            "distance": round(dist, 4) if isinstance(dist, (int, float)) else None,
        })
    return {"tool": tool, "problem": problem, "matches": matches}


def record_fix_core(tool, problem, fix, project=None, tags=None):
    cid = _hash(problem + " " + fix)
    tbl = _get_table()
    if tbl.count_rows(filter="id = '%s'" % cid):
        return {"status": "duplicate", "id": cid}
    if isinstance(tags, str):
        tags = [tags]
    tbl.add([{
        "id": cid,
        "vector": _embed(problem + " " + fix),
        "tool": tool,
        "project": project or "",
        "type": "problem_fix",
        "problem": problem,
        "fix": fix,
        "tags": json.dumps(tags or []),
        "source_path": "",
        "date_added": _now(),
    }])
    return {"status": "created", "id": cid}


def ingest_kb_core(source_path, tool=None, write_note=False):
    p = Path(source_path)
    if not p.exists():
        return {"error": "path not found: %s" % source_path}
    files = _gather_files(p)
    if not files:
        return {"ingested": 0, "skipped": 0, "errors": [],
                "note": "no supported files found under %s" % source_path}
    ingested = skipped = 0
    errors = []
    notes = []
    for f in files:
        try:
            r = _process_file(f, tool, write_note=write_note)
            if r.get("ingested"):
                ingested += r["ingested"]
            else:
                skipped += 1
            if r.get("note"):
                notes.append(r["note"])
        except Exception as exc:  # noqa: BLE001 — collect, never abort the batch
            errors.append({"file": str(f), "error": str(exc)})
    result = {"ingested": ingested, "skipped": skipped, "errors": errors}
    if notes:
        result["notes"] = notes
    return result


# --------------------------------------------------------------------------- #
# Startup auto-ingest scan
# --------------------------------------------------------------------------- #
def _scan_dir(directory, tool=None):
    """Ingest new/changed supported files under one directory. Returns the number
    of new chunks. Never raises — a broken scan is logged and returns 0."""
    try:
        directory = Path(directory)
        if not directory.exists():
            return 0
        total = 0
        for f in _gather_files(directory):
            try:
                total += _process_file(f, tool).get("ingested", 0)
            except Exception as exc:  # noqa: BLE001
                _log("ingest failed for %s: %s" % (f, exc))
        return total
    except Exception as exc:  # noqa: BLE001 — a broken scan must not block serving
        _log("scan error in %s: %s" % (directory, exc))
        return 0


def _startup_scan():
    """Ingest new/changed files on startup. Never aborts startup. Scans the
    machine-local drop folder (each file tagged by its stem) plus the vault's own
    Context/ and Reference/ folders (tagged 'vault') for system self-recall."""
    total = _scan_dir(KB_RESOURCES_DIR, tool=None)
    for d in KB_INGEST_PATHS:
        total += _scan_dir(d, tool="vault")
    if total:
        _log("startup scan ingested %d new chunks" % total)


# --------------------------------------------------------------------------- #
# Tool registry
# --------------------------------------------------------------------------- #
_RO = {"readOnlyHint": True, "openWorldHint": False}
_RW = {"readOnlyHint": False, "openWorldHint": False}


def _tool_specs():
    return [
        {"name": "query_patterns",
         "description": "Semantic search of the knowledge base for a tool, ranked by "
                        "relevance to a problem description. Returns past problem/fix "
                        "pairs and reference chunks. Call this before debugging a known tool.",
         "inputSchema": {"type": "object", "properties": {
             "tool": {"type": "string", "description": "Tool name to scope the search to."},
             "problem": {"type": "string", "description": "Describe the problem/symptom."},
             "limit": {"type": "integer", "description": "Max matches (default 5).", "default": 5}},
             "required": ["tool", "problem"]},
         "annotations": _RO, "handler": query_patterns_core},

        {"name": "list_tools",
         "description": "List every tool present in the knowledge base with its record "
                        "count. Use to see what the KB knows about.",
         "inputSchema": {"type": "object", "properties": {}, "required": []},
         "annotations": _RO, "handler": lambda: list_tools_core()},

        {"name": "get_patterns_summary",
         "description": "For one tool: record count, top tags, and the date range of "
                        "stored records. A quick profile of KB coverage for that tool.",
         "inputSchema": {"type": "object", "properties": {
             "tool": {"type": "string", "description": "Tool name to profile."}},
             "required": ["tool"]},
         "annotations": _RO, "handler": get_patterns_summary_core},

        {"name": "record_fix",
         "description": "Record a problem and its fix for a tool. Deduplicated by content "
                        "hash — recording the same problem/fix twice is a no-op. Call this "
                        "after resolving a problem worth remembering.",
         "inputSchema": {"type": "object", "properties": {
             "tool": {"type": "string", "description": "Tool the fix applies to."},
             "problem": {"type": "string", "description": "What went wrong."},
             "fix": {"type": "string", "description": "How it was resolved."},
             "project": {"type": "string", "description": "Optional project context."},
             "tags": {"type": "array", "items": {"type": "string"},
                      "description": "Optional tags."}},
             "required": ["tool", "problem", "fix"]},
         "annotations": _RW, "handler": record_fix_core},

        {"name": "ingest_kb",
         "description": "Incrementally ingest a file or directory into the knowledge base "
                        "(20 file types: docs, slides, html, epub, json/xml, csv/xlsx, "
                        "images via OCR, and plain text/pdf). Idempotent: unchanged files "
                        "and duplicate chunks are skipped. With write_note, also transposes "
                        "each non-Markdown source into a clean Obsidian note in the wiki "
                        "folder (index+note mode). Returns counts and any per-file errors.",
         "inputSchema": {"type": "object", "properties": {
             "source_path": {"type": "string", "description": "File or directory path."},
             "tool": {"type": "string", "description": "Tool name to tag ingested chunks "
                                                       "with (defaults to each file's stem)."},
             "write_note": {"type": "boolean", "description": "Also write a transposed "
                                                             "Obsidian note per source (default false)."}},
             "required": ["source_path"]},
         "annotations": _RW, "handler": ingest_kb_core},
    ]


# --------------------------------------------------------------------------- #
# Minimal MCP stdio server (JSON-RPC 2.0 over newline-delimited stdin/stdout)
# Copied from vault-mcp/server.py — identical protocol loop.
# --------------------------------------------------------------------------- #
PROTOCOL_VERSION = "2024-11-05"
SERVER_INFO = {"name": "kb", "version": "1.0.0"}


def _send(obj):
    try:
        sys.stdout.write(json.dumps(obj) + "\n")
        sys.stdout.flush()
    except BrokenPipeError:
        raise SystemExit(0)


def _serve_stdio():
    specs = _tool_specs()
    handlers = {t["name"]: t["handler"] for t in specs}
    public_tools = [{k: t[k] for k in ("name", "description", "inputSchema", "annotations")} for t in specs]

    for raw in sys.stdin:
        raw = raw.strip()
        if not raw:
            continue
        try:
            msg = json.loads(raw)
        except json.JSONDecodeError:
            continue

        mid = msg.get("id")
        method = msg.get("method")
        params = msg.get("params") or {}

        if mid is None:
            continue

        if method == "initialize":
            _send({"jsonrpc": "2.0", "id": mid, "result": {
                "protocolVersion": params.get("protocolVersion", PROTOCOL_VERSION),
                "capabilities": {"tools": {}}, "serverInfo": SERVER_INFO}})
        elif method == "tools/list":
            _send({"jsonrpc": "2.0", "id": mid, "result": {"tools": public_tools}})
        elif method == "tools/call":
            name = params.get("name")
            args = params.get("arguments") or {}
            fn = handlers.get(name)
            if fn is None:
                _send({"jsonrpc": "2.0", "id": mid,
                       "error": {"code": -32602, "message": "Unknown tool: %s" % name}})
                continue
            try:
                result = fn(**args)
                _send({"jsonrpc": "2.0", "id": mid, "result": {
                    "content": [{"type": "text", "text": json.dumps(result, ensure_ascii=False)}],
                    "structuredContent": result, "isError": False}})
            except Exception as exc:  # noqa: BLE001 — report any tool error to the client
                _send({"jsonrpc": "2.0", "id": mid, "result": {
                    "content": [{"type": "text", "text": "Error: %s" % exc}], "isError": True}})
        elif method == "ping":
            _send({"jsonrpc": "2.0", "id": mid, "result": {}})
        else:
            _send({"jsonrpc": "2.0", "id": mid,
                   "error": {"code": -32601, "message": "Method not found: %s" % method}})


if __name__ == "__main__":
    _startup_scan()
    try:
        _serve_stdio()
    except (BrokenPipeError, KeyboardInterrupt):
        pass
